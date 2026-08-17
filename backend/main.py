import os
import shutil
import tempfile
import uuid
import subprocess
from typing import Dict, Optional
from fastapi import FastAPI, File, UploadFile, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import pymupdf as fitz
import docx
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import OxmlElement
from docx.oxml.ns import qn
import pdfplumber
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
app = FastAPI(title="PDF BOX Backend API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Simple in-memory task queue
tasks: Dict[str, dict] = {}

def cleanup_temp_dir(temp_dir: str):
    try:
        shutil.rmtree(temp_dir)
    except Exception as e:
        print(f"Error cleaning up {temp_dir}: {e}")

def has_arabic(text: str) -> bool:
    return any('\u0600' <= c <= '\u06FF' for c in text)

def set_rtl_run(run):
    rPr = run._r.get_or_add_rPr()
    rtl = OxmlElement('w:rtl')
    rtl.set(qn('w:val'), '1')
    rPr.append(rtl)

def convert_pdf_to_word_task(task_id: str, input_pdf_path: str, output_docx_path: str, temp_dir: str):
    tasks[task_id]["status"] = "processing"
    
    try:
        # Assuming Windows default LibreOffice installation path
        soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        if not os.path.exists(soffice_path):
            raise Exception("LibreOffice is not installed or not found at default path.")
            
        process = subprocess.run(
            [
                soffice_path, 
                "--headless", 
                "--infilter=writer_pdf_import", 
                "--convert-to", "docx", 
                input_pdf_path, 
                "--outdir", temp_dir
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )
        
        if process.returncode != 0:
            raise Exception(f"LibreOffice conversion failed: {process.stderr}")
            
        # LibreOffice creates a file with the same basename: input.docx
        generated_docx = os.path.join(temp_dir, "input.docx")
        if os.path.exists(generated_docx):
            os.rename(generated_docx, output_docx_path)
        else:
            raise Exception("LibreOffice did not generate the output file.")
            
        tasks[task_id]["status"] = "completed"
        tasks[task_id]["download_url"] = f"/api/download/{task_id}"
        
    except Exception as e:
        tasks[task_id]["status"] = "failed"
        tasks[task_id]["error"] = str(e)
        cleanup_temp_dir(temp_dir)

@app.post("/api/convert/pdf-to-word")
async def start_pdf_to_word(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    task_id = str(uuid.uuid4())
    temp_dir = tempfile.mkdtemp()
    input_pdf_path = os.path.join(temp_dir, "input.pdf")
    output_docx_path = os.path.join(temp_dir, "output.docx")
    
    with open(input_pdf_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    tasks[task_id] = {
        "status": "pending",
        "temp_dir": temp_dir,
        "output_path": output_docx_path,
        "filename": file.filename.replace(".pdf", ".docx")
    }
    
    background_tasks.add_task(convert_pdf_to_word_task, task_id, input_pdf_path, output_docx_path, temp_dir)
    return {"task_id": task_id}

@app.get("/api/status/{task_id}")
async def get_status(task_id: str):
    if task_id not in tasks:
        return JSONResponse(status_code=404, content={"error": "Task not found"})
    return tasks[task_id]

@app.get("/api/download/{task_id}")
async def download_file(background_tasks: BackgroundTasks, task_id: str):
    if task_id not in tasks or tasks[task_id]["status"] != "completed":
        return JSONResponse(status_code=404, content={"error": "File not ready or task not found"})
        
    task_info = tasks[task_id]
    background_tasks.add_task(cleanup_temp_dir, task_info["temp_dir"])
    # Delete from dict to prevent memory leak
    del tasks[task_id]
    
    return FileResponse(
        task_info["output_path"],
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=task_info["filename"]
    )

def convert_pdf_to_ppt_task(task_id: str, input_pdf_path: str, output_pptx_path: str, temp_dir: str):
    tasks[task_id]["status"] = "processing"
    
    try:
        prs = Presentation()
        doc = fitz.open(input_pdf_path)
        
        # Blank slide layout
        blank_slide_layout = prs.slide_layouts[6] 
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            blocks = page.get_text("dict")["blocks"]
            blocks.sort(key=lambda b: (b["bbox"][1], b["bbox"][0]))
            
            for block in blocks:
                if block["type"] == 0:
                    x0, y0, x1, y1 = block["bbox"]
                    
                    # Convert fitz coordinates (points) to python-pptx coordinates
                    # 1 inch = 72 points = 914400 EMU. Therefore 1 point = 12700 EMU
                    left = int(x0 * 12700)
                    top = int(y0 * 12700)
                    width = int((x1 - x0) * 12700)
                    height = int((y1 - y0) * 12700)
                    
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    lines = block["lines"]
                    lines.sort(key=lambda l: l["bbox"][1])
                    
                    first_line = True
                    for line in lines:
                        words = []
                        for span in line["spans"]:
                            text = span["text"].strip()
                            if text:
                                words.append(text)
                                
                        if not words:
                            continue
                            
                        line_text = " ".join(words)
                        is_arabic_line = has_arabic(line_text)
                        
                        if first_line:
                            p = tf.paragraphs[0]
                            first_line = False
                        else:
                            p = tf.add_paragraph()
                            
                        if is_arabic_line:
                            p.alignment = PP_ALIGN.RIGHT
                        else:
                            p.alignment = PP_ALIGN.LEFT
                            
                        run = p.add_run()
                        run.text = line_text
                        
        prs.save(output_pptx_path)
        doc.close()
        
        tasks[task_id]["status"] = "completed"
        tasks[task_id]["download_url"] = f"/api/download/{task_id}"
        
    except Exception as e:
        tasks[task_id]["status"] = "failed"
        tasks[task_id]["error"] = str(e)
        cleanup_temp_dir(temp_dir)


@app.post("/api/convert/pdf-to-ppt")
async def start_pdf_to_ppt(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    task_id = str(uuid.uuid4())
    temp_dir = tempfile.mkdtemp()
    input_pdf_path = os.path.join(temp_dir, "input.pdf")
    output_pptx_path = os.path.join(temp_dir, "output.pptx")
    
    with open(input_pdf_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    tasks[task_id] = {
        "status": "pending",
        "temp_dir": temp_dir,
        "output_path": output_pptx_path,
        "filename": file.filename.replace(".pdf", ".pptx")
    }
    
    background_tasks.add_task(convert_pdf_to_ppt_task, task_id, input_pdf_path, output_pptx_path, temp_dir)
    return {"task_id": task_id}

def office_to_pdf_task(task_id: str, input_path: str, temp_dir: str):
    tasks[task_id]["status"] = "processing"
    
    try:
        # Assuming Windows default LibreOffice installation path
        soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        if not os.path.exists(soffice_path):
            raise Exception("LibreOffice is not installed or not found at default path.")
            
        process = subprocess.run(
            [soffice_path, "--headless", "--convert-to", "pdf", input_path, "--outdir", temp_dir],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )
        
        if process.returncode != 0:
            raise Exception(f"LibreOffice conversion failed: {process.stderr}")
            
        tasks[task_id]["status"] = "completed"
        tasks[task_id]["download_url"] = f"/api/download/{task_id}"
        
    except Exception as e:
        tasks[task_id]["status"] = "failed"
        tasks[task_id]["error"] = str(e)
        cleanup_temp_dir(temp_dir)


@app.post("/api/convert/office-to-pdf")
async def start_office_to_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    task_id = str(uuid.uuid4())
    temp_dir = tempfile.mkdtemp()
    
    # Preserve original extension
    ext = os.path.splitext(file.filename)[1].lower()
    input_path = os.path.join(temp_dir, f"input{ext}")
    output_pdf_path = os.path.join(temp_dir, "input.pdf") # LibreOffice creates a PDF with the same basename
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    tasks[task_id] = {
        "status": "pending",
        "temp_dir": temp_dir,
        "output_path": output_pdf_path,
        "filename": os.path.splitext(file.filename)[0] + ".pdf"
    }
    
    background_tasks.add_task(office_to_pdf_task, task_id, input_path, temp_dir)
    return {"task_id": task_id}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)

def convert_pdf_to_excel_task(task_id: str, input_pdf_path: str, output_xlsx_path: str, temp_dir: str):
    tasks[task_id]["status"] = "processing"
    
    try:
        wb = Workbook()
        ws = wb.active
        ws.title = "Page 1"
        
        has_arabic_in_doc = False
        
        with pdfplumber.open(input_pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                if page_num > 0:
                    ws = wb.create_sheet(title=f"Page {page_num + 1}")
                    
                tables = page.extract_tables()
                if not tables:
                    ws.append(["No tables found on this page"])
                    continue
                    
                for table in tables:
                    for row in table:
                        clean_row = []
                        for cell in row:
                            if cell:
                                cell_text = str(cell).replace('\n', ' ')
                                if has_arabic(cell_text):
                                    has_arabic_in_doc = True
                                clean_row.append(cell_text)
                            else:
                                clean_row.append("")
                        ws.append(clean_row)
                    ws.append([]) # Empty row between tables
                
                # If Arabic is detected on this sheet, flip the sheet to RTL
                if has_arabic_in_doc:
                    ws.sheet_view.rightToLeft = True
                    has_arabic_in_doc = False # Reset for next page
                    
        wb.save(output_xlsx_path)
        
        tasks[task_id]["status"] = "completed"
        tasks[task_id]["download_url"] = f"/api/download/{task_id}"
        
    except Exception as e:
        tasks[task_id]["status"] = "failed"
        tasks[task_id]["error"] = str(e)
        cleanup_temp_dir(temp_dir)

@app.post("/api/convert/pdf-to-excel")
async def start_pdf_to_excel(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    task_id = str(uuid.uuid4())
    temp_dir = tempfile.mkdtemp()
    input_pdf_path = os.path.join(temp_dir, "input.pdf")
    output_xlsx_path = os.path.join(temp_dir, "output.xlsx")
    
    with open(input_pdf_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    tasks[task_id] = {
        "status": "pending",
        "temp_dir": temp_dir,
        "output_path": output_xlsx_path,
        "filename": file.filename.replace(".pdf", ".xlsx")
    }
    
    background_tasks.add_task(convert_pdf_to_excel_task, task_id, input_pdf_path, output_xlsx_path, temp_dir)
    return {"task_id": task_id}
